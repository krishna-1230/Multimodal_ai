from flask import Flask, request, send_file, jsonify, send_from_directory
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from io import BytesIO
import os
import uuid
import logging
from flask_cors import CORS
import requests

app = Flask(__name__)
CORS(app)

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


@app.route('/convert-pptx', methods=['POST'])
def convert_to_pptx():
    try:
        data = request.get_json(force=True)
    except Exception:
        return jsonify({'error': 'Invalid JSON body'}), 400

    if not data:
        return jsonify({'error': 'No JSON data provided'}), 400

    prs = Presentation()

    # Global slide size support: optional values '16:9' or '4:3' or custom dict with width/height in inches
    slide_size = data.get('slide_size')
    try:
        if isinstance(slide_size, str):
            if slide_size == '16:9':
                prs.slide_width = Inches(13.3333)
                prs.slide_height = Inches(7.5)
            elif slide_size == '4:3':
                prs.slide_width = Inches(10)
                prs.slide_height = Inches(7.5)
        elif isinstance(slide_size, dict):
            w = float(slide_size.get('width', 10))
            h = float(slide_size.get('height', 7.5))
            prs.slide_width = Inches(w)
            prs.slide_height = Inches(h)
    except Exception:
        # ignore invalid slide_size, continue with defaults
        pass

    for slide_data in data.get('slides', []):
        # Use a safe layout index
        layout_index = slide_data.get('layout_index', 1)
        layout_index = int(layout_index) if isinstance(layout_index, (int, str)) else 1
        layout_index = max(0, min(layout_index, len(prs.slide_layouts) - 1))
        slide_layout = prs.slide_layouts[layout_index]
        slide = prs.slides.add_slide(slide_layout)

        # Add title if present
        title_text = slide_data.get('title')
        if title_text and slide.shapes.title:
            slide.shapes.title.text = str(title_text)
            # Title style
            title_style = slide_data.get('title_style', {})
            try:
                title_tf = slide.shapes.title.text_frame
                for p in title_tf.paragraphs:
                    for run in p.runs:
                        if title_style.get('font_size'):
                            run.font.size = Pt(float(title_style.get('font_size')))
                        if title_style.get('font_name'):
                            run.font.name = title_style.get('font_name')
                        color = title_style.get('color')
                        if color and isinstance(color, str) and color.startswith('#') and len(color) == 7:
                            try:
                                r = int(color[1:3], 16)
                                g = int(color[3:5], 16)
                                b = int(color[5:7], 16)
                                run.font.color.rgb = RGBColor(r, g, b)
                            except Exception:
                                pass
            except Exception:
                pass

        # Add content to body placeholder
        body_shape = None
        for shape in slide.shapes:
            if getattr(shape, 'has_text_frame', False) and shape != slide.shapes.title:
                body_shape = shape
                break

        if body_shape:
            tf = body_shape.text_frame
            # Clear any existing paragraphs
            tf.clear()
            for item in slide_data.get('content', []):
                item_type = item.get('type')
                item_style = item.get('style', {})

                if item_type == 'paragraph':
                    p = tf.add_paragraph()
                    run = p.add_run()
                    run.text = str(item.get('text', ''))
                    if item_style.get('font_size'):
                        run.font.size = Pt(float(item_style.get('font_size')))
                    if item_style.get('font_name'):
                        run.font.name = item_style.get('font_name')
                    color = item_style.get('color')
                    if color and isinstance(color, str) and color.startswith('#') and len(color) == 7:
                        try:
                            r = int(color[1:3], 16)
                            g = int(color[3:5], 16)
                            b = int(color[5:7], 16)
                            run.font.color.rgb = RGBColor(r, g, b)
                        except Exception:
                            pass
                elif item_type == 'list':
                    for li_item in item.get('items', []):
                        p = tf.add_paragraph()
                        run = p.add_run()
                        run.text = str(li_item)
                        p.level = int(item.get('level', 1))
                elif item_type in ('subtitle', 'heading'):
                    left = Inches(item.get('left', 1))
                    top = Inches(item.get('top', 1.5))
                    width = Inches(item.get('width', 8))
                    height = Inches(item.get('height', 0.5))
                    box = slide.shapes.add_textbox(left, top, width, height)
                    tf_box = box.text_frame
                    p_box = tf_box.add_paragraph()
                    run = p_box.add_run()
                    run.text = str(item.get('text', ''))
                    if item_style.get('font_size'):
                        run.font.size = Pt(float(item_style.get('font_size')))
                    if item_style.get('font_name'):
                        run.font.name = item_style.get('font_name')
                    if item_style.get('bold'):
                        run.font.bold = True
                    color = item_style.get('color')
                    if color and isinstance(color, str) and color.startswith('#') and len(color) == 7:
                        try:
                            r = int(color[1:3], 16)
                            g = int(color[3:5], 16)
                            b = int(color[5:7], 16)
                            run.font.color.rgb = RGBColor(r, g, b)
                        except Exception:
                            pass
                elif item_type == 'table':
                    rows = int(item.get('rows', 1))
                    cols = int(item.get('cols', 1))
                    left = Inches(item.get('left', 1))
                    top = Inches(item.get('top', 1))
                    width = Inches(item.get('width', 8))
                    height = Inches(item.get('height', 1))
                    shape = slide.shapes.add_table(rows, cols, left, top, width, height)
                    table = shape.table

                    for r_idx, row_data in enumerate(item.get('data', [])):
                        for c_idx, cell_text in enumerate(row_data):
                            if r_idx < rows and c_idx < cols:
                                cell = table.cell(r_idx, c_idx)
                                cell.text = str(cell_text)

        # Handle images separately as they are not part of text frame
        # Slide-level settings: background color
        bg = slide_data.get('background')
        if bg and isinstance(bg, dict):
            color = bg.get('color')
            if color and isinstance(color, str) and color.startswith('#') and len(color) == 7:
                try:
                    slide.background.fill.solid()
                    r = int(color[1:3], 16)
                    g = int(color[3:5], 16)
                    b = int(color[5:7], 16)
                    slide.background.fill.fore_color.rgb = RGBColor(r, g, b)
                except Exception:
                    pass

        for item in slide_data.get('content', []):
            if item.get('type') == 'image':
                left = Inches(item.get('left', 1))
                top = Inches(item.get('top', 1))
                width = Inches(item.get('width', 6))
                height = Inches(item.get('height', 4.5))
                # Support either local path or remote URL
                img_path = item.get('path')
                try:
                    if isinstance(img_path, str) and img_path.startswith(('http://', 'https://')):
                        resp = requests.get(img_path, timeout=10)
                        resp.raise_for_status()
                        img_bytes = BytesIO(resp.content)
                        slide.shapes.add_picture(img_bytes, left, top, width=width, height=height)
                    else:
                        slide.shapes.add_picture(img_path, left, top, width=width, height=height)
                except Exception as e:
                    logger.warning('Error adding image to slide: %s', e)

    # Ensure the directory for generated PPTX files exists
    output_dir = 'output'
    os.makedirs(output_dir, exist_ok=True)

    # Generate a unique filename
    unique_filename = f"{uuid.uuid4()}.pptx"
    output_path = os.path.join(output_dir, unique_filename)

    # Save the presentation to a file
    prs.save(output_path)

    # Construct the response JSON
    file_url = request.url_root.rstrip('/') + f'/{output_dir}/{unique_filename}'

    return jsonify({
        "success": True,
        "file": {
            "name": unique_filename,
            "path": os.path.abspath(output_path),
            "downloadUrl": file_url
        }
    })


@app.route('/output/<path:filename>')
def download_pptx(filename):
    return send_from_directory('output', filename, as_attachment=True)


if __name__ == '__main__':
    app.run(debug=True, port=5002)